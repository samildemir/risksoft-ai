import tempfile
import docx
import subprocess
from pptx import Presentation
from io import BytesIO
from typing import List
import logging
from utils.helper import replaceName
import re
from nltk.corpus import stopwords
from utils.s3Handler import S3Handler
from utils.helper import detect_file_type
import logging
from io import BytesIO
from langchain.docstore.document import Document as LangchainDocument
from langchain.text_splitter import CharacterTextSplitter
from docx import Document as DocxDocument
from pptx import Presentation
import subprocess
from typing import List, Dict, Optional
from models.schemas import DocumentSource
import fitz
from constants.config import GPT_4o, CHUNK_SIZE_LIMIT, MAX_CHUNK_OVERLAP



def process_pdf(file_content):
    pdf_document = fitz.open(stream=BytesIO(file_content), filetype="pdf")
    pages = []
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        pages.append(page.get_text())
    return pages

def process_doc(file_content):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as temp_file:
        temp_file.write(file_content)
        temp_file_path = temp_file.name

    result = subprocess.run(['antiword', temp_file_path], capture_output=True, text=True)
    return [result.stdout]

def process_docx(file_content):
    doc = docx.Document(BytesIO(file_content))
    return ["\n\n".join([para.text for para in doc.paragraphs])]

def process_pptx(file_content):
    ppt = Presentation(BytesIO(file_content))
    content = []
    for slide in ppt.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_content.append(shape.text)
        content.append("\n".join(slide_content))
    return content

def download_and_order_files(path_list: List[DocumentSource], account_id, raw_text: bool = False):
    try:
        client = S3Handler()
        doc_list = []
        for data in path_list:
            file_content = client.s3_get_object(data.path)
            file_type = detect_file_type(file_content)

            if file_type == "application/pdf":
                pages = process_pdf(file_content)
            elif file_type == "application/msword":
                pages = process_doc(file_content)
            elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                pages = process_docx(file_content)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                pages = process_pptx(file_content)
            else:
                continue

            if raw_text:  # Check if raw_text is True
                return "\n".join(pages)  # Return raw text as a single string

            for i, page in enumerate(pages):
                chunks = get_text_chunks(page)
                for j, chunk in enumerate(chunks):
                    doc_list.append(LangchainDocument(
                        page_content=f"Document: {data.name}, Page: {i + 1}\n{chunk}",
                        metadata={
                            "account_id": account_id,
                            "source": "document",
                            "file_name": data.name,
                            "description": replaceName(data.name),
                            "page_number": i + 1,
                            "chunk_id": j,
                            "total_chunks": len(chunks)
                        }
                    ))
        return doc_list
    except Exception as e:
        logging.error(f"Error downloading and ordering files: {e}")
        return []


def preprocess_text(text):
    try:
        # Remove non-alphabetic characters and normalize whitespace
        text = re.sub(r'[^a-zA-ZğüşöçıİĞÜŞÖÇ\s]+', ' ', text).strip()
        
        # Convert to lowercase
        text = text.lower()
        
        # Remove stop words
        stop_words = set(stopwords.words('turkish'))
        filtered_words = [word for word in text.split() if word not in stop_words]
        
        return ' '.join(filtered_words)
    except Exception as e:
        print(f"An error occurred: {e}")
        return ""

def get_text_chunks(text: str) -> List[str]:
    text = preprocess_text(text)
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=CHUNK_SIZE_LIMIT,
        chunk_overlap=MAX_CHUNK_OVERLAP,
        length_function=len,
    )
    chunks = text_splitter.split_text(text)
    return chunks